# -*- coding: utf-8 -*-
# ---------------------------------------------------------------------
# Copyright (c) 2026 Qualcomm Technologies, Inc. and/or its subsidiaries.
# SPDX-License-Identifier: BSD-3-Clause
# ---------------------------------------------------------------------
"""
pptx_base.py — 轻量可编辑 PPTX 基础库（通用深色鎏金典雅风）

目标：用 python-pptx 原生对象复用参考 PPT 中的通用设计语言：
深棕黑底、朱红/鎏金/青绿强调线、仪式感边框、目录卡片、时间轴、KPI、三列卡片、左右分栏、图文页。

所有主要元素均为 PowerPoint 可编辑对象：文本框、矩形、圆形、线条、表格与可替换图片。
禁止把整页渲染成一张大图。可选图片只用于封面/插图/背景氛围，叠加文字与卡片仍保持可编辑。
"""

from datetime import datetime
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ─────────────────────────────────────────────────────────────────────────────
# 画布与主题常量
# ─────────────────────────────────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 13.333, 7.5
SAFE_L, SAFE_R = 0.65, 12.7
FONT_TITLE = "KaiTi"
FONT_BODY = "Microsoft YaHei"
FONT_SONG = "SimSun"

BG_DARK = RGBColor(0x18, 0x0B, 0x06)
BG_DEEP = RGBColor(0x20, 0x10, 0x08)
BG_CARD = RGBColor(0x2B, 0x16, 0x0B)
BG_CARD2 = RGBColor(0x37, 0x1E, 0x0F)
BG_INK = RGBColor(0x0F, 0x07, 0x04)
BORDER_DIM = RGBColor(0x56, 0x32, 0x18)
BORDER_SOFT = RGBColor(0x39, 0x22, 0x12)

ACCENT_RED = RGBColor(0xC6, 0x3A, 0x2B)
ACCENT_YEL = RGBColor(0xD6, 0xA6, 0x16)
ACCENT_GRN = RGBColor(0x35, 0x91, 0x68)
ACCENT_CREAM = RGBColor(0xC9, 0xB5, 0x8A)
ACCENT_ORANGE = RGBColor(0xC7, 0x63, 0x24)
ACCENT_BLUE = ACCENT_YEL
ACCENT_JADE = ACCENT_GRN
ACCENT_BRONZE = RGBColor(0x8A, 0x5A, 0x22)

TEXT_WHITE = RGBColor(0xF4, 0xE7, 0xC8)
TEXT_LIGHT = RGBColor(0xD9, 0xBD, 0x88)
TEXT_DIM = RGBColor(0x9A, 0x7A, 0x54)
TITLE_BLUE = RGBColor(0xD8, 0x90, 0x24)
COLOR_POS = ACCENT_GRN
COLOR_NEG = ACCENT_RED
BRAND_COLORS = [ACCENT_RED, ACCENT_YEL, ACCENT_GRN, ACCENT_CREAM]


def _visual_len(text) -> float:
    """估算字符串视觉长度：中文/全角按 1，英文/数字/空格按更窄比例。"""
    total = 0.0
    for ch in str(text):
        if ch == " ":
            total += 0.32
        elif ord(ch) < 128:
            total += 0.55
        else:
            total += 1.0
    return total


def _fit_font_size(text, box_w, base_size, min_size=8.0, density=2.25):
    """按文本长度与文本框宽度粗略压缩字号，避免长标题/英文串溢出。"""
    vlen = max(_visual_len(text), 1.0)
    capacity = max(box_w * density, 1.0)
    if vlen <= capacity:
        return base_size
    return max(min_size, base_size * capacity / vlen)

# ─────────────────────────────────────────────────────────────────────────────
# 基础对象
# ─────────────────────────────────────────────────────────────────────────────
def init_presentation(export_dir: str):
    """创建 16:9 PPTX，并确保输出目录存在。返回 (prs, blank_layout)。"""
    os.makedirs(export_dir, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs, prs.slide_layouts[6]


def _set_transparency(shape, transparency: float):
    """设置形状透明度；0=不透明，1=全透明。旧版 python-pptx 不支持时自动忽略。"""
    try:
        shape.fill.transparency = max(0, min(1, transparency))
    except Exception:
        pass


def add_shape(slide, shape_type, x, y, w, h, fill_color=None, line_color=None,
              line_width=0.8, transparency=0.0, radius=True):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        _set_transparency(shape, transparency)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    if radius and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=0.0, transparency=0.0):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill_color, line_color, line_width, transparency, False)


def add_rounded_rect(slide, x, y, w, h, fill_color=BG_CARD, line_color=None, line_width=0.8, transparency=0.0):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill_color, line_color, line_width, transparency, True)


def add_circle(slide, x, y, d, fill_color=None, line_color=None, line_width=0.8, transparency=0.0):
    return add_shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill_color, line_color, line_width, transparency, False)


def add_line(slide, x1, y1, x2, y2, color=BORDER_DIM, width=0.8):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_text(slide, text, x, y, w, h, font_size=16, bold=False, color=TEXT_WHITE,
             align=PP_ALIGN.LEFT, font_name=FONT_BODY, valign=MSO_ANCHOR.TOP,
             auto_fit=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    final_size = _fit_font_size(text, w, font_size) if auto_fit else font_size
    run.font.size = Pt(final_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return box


def add_multiline_text(slide, lines, x, y, w, h, font_size=13, color=TEXT_LIGHT,
                       align=PP_ALIGN.LEFT, bold=False, font_name=FONT_BODY, line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = str(line)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color
    return box


def add_bg_picture(slide, image_path: str, overlay_alpha: float = 0.42):
    """添加全页图片背景和深色遮罩。注意：图片可编辑/可替换，但文字与卡片仍是原生对象。"""
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_DARK, transparency=overlay_alpha)
        return True
    return False

# ─────────────────────────────────────────────────────────────────────────────
# 深色鎏金底纹、标题、页脚
# ─────────────────────────────────────────────────────────────────────────────
def add_tang_texture(slide, with_frame=True, top_accent=None, bottom_accent=None):
    """深色鎏金背景：底色、低对比细网格、内收暗纹、边框，均为可编辑形状。"""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_DARK)
    # 氛围圆必须完全落在画布内，且颜色接近背景，避免编辑视图中出现"越界大黑圆"。
    add_circle(slide, 4.95, 0.15, 3.05, RGBColor(0x22, 0x10, 0x07), line_color=None)
    add_circle(slide, 8.25, 4.55, 2.75, RGBColor(0x20, 0x0F, 0x07), line_color=None)
    # 网格只保留在内容安全区内，低对比，避免抢正文层级。
    grid_color = RGBColor(0x2A, 0x18, 0x0D)
    for x in [1.45, 2.95, 4.45, 5.95, 7.45, 8.95, 10.45, 11.95]:
        add_rect(slide, x, 0.78, 0.006, 5.92, grid_color)
    for y in [1.28, 2.36, 3.44, 4.52, 5.60, 6.68]:
        add_rect(slide, 0.62, y, 12.08, 0.006, grid_color)
    if top_accent:
        add_rect(slide, 0, 0.05, SLIDE_W, 0.045, top_accent)
    if bottom_accent:
        add_rect(slide, 0, 7.40, SLIDE_W, 0.045, bottom_accent)
    if with_frame:
        add_line(slide, 0.65, 1.12, 12.75, 1.12, BORDER_DIM, 0.8)
        add_line(slide, 0.65, 6.88, 12.75, 6.88, BORDER_DIM, 0.8)


def add_corner_frame(slide, color=ACCENT_YEL, inset=0.62, length=0.92, width=1.1):
    for x, y, sx, sy in [(inset, 0.55, 1, 1), (SLIDE_W - inset, 0.55, -1, 1),
                         (inset, SLIDE_H - 0.55, 1, -1), (SLIDE_W - inset, SLIDE_H - 0.55, -1, -1)]:
        add_line(slide, x, y, x + sx * length, y, color, width)
        add_line(slide, x, y, x, y + sy * length, color, width)


def add_chapter_tag(slide, text, color=ACCENT_RED, x=0.68, y=0.28):
    tag_w = max(1.75, min(3.1, _visual_len(text) * 0.12 + 0.55))
    add_rounded_rect(slide, x, y, tag_w, 0.30, color, line_color=color, line_width=0)
    add_text(slide, text, x + 0.05, y + 0.055, tag_w - 0.1, 0.16, 8.2, True, TEXT_WHITE, PP_ALIGN.CENTER)
    return tag_w


def add_imperial_title(slide, title, subtitle="", chapter_tag="", accent=ACCENT_RED):
    """内容页标题：左竖线、书法标题、副标题、右上标签。与 add_chapter_tag 默认位置错开。"""
    add_rect(slide, 0.65, 0.66, 0.045, 0.42, accent)
    title_fs = _fit_font_size(title, 10.9 if not chapter_tag else 9.7, 26, min_size=18, density=2.05)
    add_text(slide, title, 0.88, 0.56, 10.9 if not chapter_tag else 9.7, 0.42,
             title_fs, False, TEXT_WHITE, font_name=FONT_TITLE, auto_fit=False)
    if subtitle:
        add_text(slide, subtitle, 0.88, 1.06, 9.6, 0.20, 10.2, False, TEXT_LIGHT)
    add_line(slide, 0.65, 1.34, 12.65, 1.34, BORDER_DIM, 0.7)
    add_line(slide, 0.65, 1.34, 4.5, 1.34, ACCENT_YEL, 0.9)
    if chapter_tag:
        add_rounded_rect(slide, 10.85, 0.65, 1.85, 0.30, BG_CARD2, BORDER_DIM, 0.8)
        add_text(slide, chapter_tag, 10.9, 0.70, 1.75, 0.20, 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)


def add_imperial_footer(slide, section_text, num=None, total=None):
    add_line(slide, 0.65, 6.88, 12.65, 6.88, BORDER_DIM, 0.8)
    if section_text:
        add_text(slide, section_text, 0.65, 7.03, 5.8, 0.20, 8, False, TEXT_DIM)
    if num is not None and total is not None:
        add_text(slide, f"{num:02d} / {total}", 12.05, 7.03, 0.75, 0.20, 8, False, TEXT_DIM, PP_ALIGN.RIGHT)


def add_data_source(slide, source_text):
    add_text(slide, f"来源：{source_text}", 0.65, 6.78, 6, 0.18, 7.5, False, TEXT_DIM)


def add_imperial_takeaway(slide, text, y=6.46, color=ACCENT_YEL):
    add_rounded_rect(slide, 0.78, y, 11.85, 0.30, BG_CARD, line_color=None, transparency=0.10)
    add_text(slide, text, 1.05, y + 0.04, 11.3, 0.20, 10, True, color, PP_ALIGN.CENTER, FONT_SONG)

# 语义化别名：保留旧函数名兼容历史脚本，同时让新脚本不再绑定任何固定题材语义。
def add_gilded_texture(slide, with_frame=True, top_accent=None, bottom_accent=None):
    return add_tang_texture(slide, with_frame=with_frame, top_accent=top_accent, bottom_accent=bottom_accent)


def add_section_tag(slide, text, color=ACCENT_RED, x=0.68, y=0.28):
    return add_chapter_tag(slide, text, color=color, x=x, y=y)


def add_deck_title(slide, title, subtitle="", section_note="", accent=ACCENT_RED):
    return add_imperial_title(slide, title, subtitle=subtitle, chapter_tag=section_note, accent=accent)


def add_deck_footer(slide, section_text, num=None, total=None):
    return add_imperial_footer(slide, section_text, num=num, total=total)


def add_deck_takeaway(slide, text, y=6.46, color=ACCENT_YEL):
    return add_imperial_takeaway(slide, text, y=y, color=color)

# ─────────────────────────────────────────────────────────────────────────────
# 结构页
# ─────────────────────────────────────────────────────────────────────────────
def add_cover_slide(prs, blank, title, subtitle="", author="", date="",
                    tag_text="专题汇报 · STRATEGIC BRIEF", year_range="适用范围 / 时间周期", bg_image="",
                    page_text="01 / --"):
    slide = prs.slides.add_slide(blank)
    if not add_bg_picture(slide, bg_image, 0.45):
        add_tang_texture(slide, with_frame=False)
    add_corner_frame(slide, ACCENT_YEL, inset=0.72, length=0.9)
    add_rounded_rect(slide, 5.25, 1.42, 2.85, 0.34, ACCENT_RED, line_color=None)
    add_text(slide, tag_text, 5.32, 1.48, 2.71, 0.20, 10.5, False, TEXT_WHITE, PP_ALIGN.CENTER)
    add_text(slide, title, 2.2, 2.65, 8.9, 0.88, 54, False, TITLE_BLUE, PP_ALIGN.CENTER, FONT_TITLE)
    add_line(slide, 3.0, 3.62, 6.0, 3.62, ACCENT_YEL, 0.9)
    add_line(slide, 7.35, 3.62, 10.35, 3.62, ACCENT_YEL, 0.9)
    add_text(slide, "◆", 6.45, 3.49, 0.45, 0.25, 13, False, ACCENT_YEL, PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 2.1, 4.00, 9.2, 0.38, 18, True, TEXT_WHITE, PP_ALIGN.CENTER)
    if year_range:
        add_rounded_rect(slide, 5.25, 4.58, 2.85, 0.42, BG_DEEP, ACCENT_YEL, 1.0)
        add_text(slide, year_range, 5.35, 4.66, 2.65, 0.24, 12, True, ACCENT_YEL, PP_ALIGN.CENTER)
    if author or date:
        add_text(slide, f"{author} · {date}" if date else author, 2.0, 6.45, 9.3, 0.24, 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    for i, c in enumerate([ACCENT_RED, ACCENT_YEL, ACCENT_GRN, ACCENT_CREAM, ACCENT_YEL]):
        add_circle(slide, 5.95 + i * 0.22, 6.92, 0.055, c, c)
    add_text(slide, page_text, 6.15, 7.18, 1.0, 0.18, 8, False, TEXT_DIM, PP_ALIGN.CENTER)
    return slide


def add_ending_slide(prs, blank, closing_text="感谢观看", contact="", sub_text="", bg_image="",
                     tag_text="SUMMARY · NEXT STEP", year_range="适用范围 / 时间周期",
                     page_text="-- / --"):
    slide = prs.slides.add_slide(blank)
    if not add_bg_picture(slide, bg_image, 0.50):
        add_tang_texture(slide, with_frame=False, top_accent=ACCENT_YEL, bottom_accent=ACCENT_YEL)
    add_corner_frame(slide, ACCENT_YEL, inset=0.72, length=0.9)
    add_rounded_rect(slide, 4.85, 1.38, 3.65, 0.34, ACCENT_YEL, line_color=None)
    add_text(slide, tag_text, 4.95, 1.45, 3.45, 0.20, 9.5, False, BG_DARK, PP_ALIGN.CENTER)
    add_circle(slide, 4.45, 1.85, 4.25, BG_INK, BORDER_DIM, 0.8, transparency=0.20)
    title_size = _fit_font_size(closing_text, 8.9, 54, min_size=34, density=1.8)
    add_text(slide, closing_text, 2.2, 2.65, 8.9, 0.88, title_size, False, TITLE_BLUE, PP_ALIGN.CENTER, FONT_TITLE, auto_fit=False)
    add_line(slide, 3.0, 3.62, 6.0, 3.62, ACCENT_YEL, 0.9)
    add_line(slide, 7.35, 3.62, 10.35, 3.62, ACCENT_YEL, 0.9)
    add_text(slide, "◆", 6.45, 3.49, 0.45, 0.25, 13, False, ACCENT_YEL, PP_ALIGN.CENTER)
    if sub_text:
        add_multiline_text(slide, sub_text.split("\n"), 3.2, 3.92, 6.9, 0.78, 16, TEXT_WHITE, PP_ALIGN.CENTER, False, FONT_TITLE)
    if year_range:
        add_rounded_rect(slide, 5.25, 4.82, 2.85, 0.42, BG_DEEP, ACCENT_YEL, 1.0)
        add_text(slide, year_range, 5.35, 4.90, 2.65, 0.24, 12, True, ACCENT_YEL, PP_ALIGN.CENTER)
    if contact:
        add_text(slide, contact, 2.1, 6.45, 9.2, 0.24, 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    add_text(slide, page_text, 6.15, 7.18, 1.0, 0.18, 8, False, TEXT_DIM, PP_ALIGN.CENTER)
    return slide


def add_toc_slide(prs, blank, toc_items, total_pages, title="目录", subtitle="Contents Overview", toc_details=None):
    slide = prs.slides.add_slide(blank)
    add_tang_texture(slide, with_frame=False)
    add_rect(slide, 0.66, 0.58, 0.045, 0.45, ACCENT_RED)
    toc_title = f"{title} — {subtitle}"
    toc_font = FONT_BODY if any(ord(ch) < 128 and ch.isalpha() for ch in toc_title) else FONT_TITLE
    add_text(slide, toc_title, 0.88, 0.52, 8.2, 0.42, 21, True, TEXT_WHITE, font_name=toc_font)
    add_line(slide, 0.65, 1.17, 12.65, 1.17, BORDER_DIM, 0.8)
    add_line(slide, 0.65, 1.17, 3.55, 1.17, ACCENT_YEL, 1.0)
    add_rounded_rect(slide, 10.92, 0.68, 1.8, 0.30, BG_CARD2, BORDER_DIM, 0.8)
    add_text(slide, f"{title} · {total_pages}页", 10.98, 0.73, 1.68, 0.18, 8.5, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    positions = [(0.68, 1.42), (6.85, 1.42), (0.68, 4.25), (6.85, 4.25)]
    colors = [ACCENT_RED, ACCENT_YEL, ACCENT_GRN, ACCENT_CREAM]
    card_w, card_h = 5.7, 2.52
    details = toc_details or []
    for i, item in enumerate(toc_items[:4]):
        x, y = positions[i]
        d = details[i] if i < len(details) else {}
        c = d.get("color", colors[i % 4])
        add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD, c, 1.0)
        add_rect(slide, x, y, 0.055, card_h, c)
        tag = d.get("tag", f"第{i + 1}章")
        tag_w = max(0.86, min(1.25, _visual_len(tag) * 0.12 + 0.25))
        add_rounded_rect(slide, x + 0.22, y + 0.18, tag_w, 0.28, c, line_color=None)
        add_text(slide, tag, x + 0.26, y + 0.23, tag_w - 0.08, 0.16, 8.0, True, TEXT_WHITE, PP_ALIGN.CENTER)
        add_text(slide, item, x + tag_w + 0.48, y + 0.15, card_w - tag_w - 0.75, 0.35, 18, False, ACCENT_YEL, font_name=FONT_TITLE)
        if d.get("desc"):
            add_text(slide, d["desc"], x + 0.25, y + 0.72, 4.8, 0.24, 10, False, TEXT_LIGHT)
        add_line(slide, x + 0.25, y + 0.88, x + card_w - 0.35, y + 0.88, BORDER_DIM, 0.6)
        pages = d.get("pages", [])
        for j, p in enumerate(pages[:3]):
            label, desc = p if isinstance(p, (tuple, list)) and len(p) >= 2 else (str(p), "")
            py = y + 1.08 + j * 0.38
            add_circle(slide, x + 0.33, py + 0.06, 0.06, c, c)
            add_text(slide, label, x + 0.52, py, 2.6, 0.22, 11, True, TEXT_WHITE)
            add_text(slide, desc, x + 3.08, py, 2.15, 0.22, 9.5, False, TEXT_DIM, PP_ALIGN.RIGHT)
        if d.get("page_range"):
            add_text(slide, d["page_range"], x + 0.25, y + card_h - 0.28, 1.5, 0.18, 8.5, False, c)
    add_imperial_footer(slide, "目录", 2, total_pages)
    return slide


def add_section_slide(prs, blank, section_num, section_title, section_subtitle=""):
    slide = prs.slides.add_slide(blank)
    add_tang_texture(slide, with_frame=False)
    add_rect(slide, 0, 0, 0.07, SLIDE_H / 3, ACCENT_RED)
    add_rect(slide, 0, SLIDE_H / 3, 0.07, SLIDE_H / 3, ACCENT_YEL)
    add_rect(slide, 0, SLIDE_H * 2 / 3, 0.07, SLIDE_H / 3, ACCENT_GRN)
    add_text(slide, section_num, 0.95, 1.10, 4.5, 2.0, 132, True, RGBColor(0x31, 0x19, 0x0B))
    add_text(slide, section_title, 1.0, 4.25, 9.8, 0.72, 42, False, TEXT_WHITE, font_name=FONT_TITLE)
    if section_subtitle:
        add_text(slide, section_subtitle, 1.05, 5.07, 8.8, 0.34, 16, False, TEXT_DIM)
    add_line(slide, 1.0, 5.72, 5.2, 5.72, ACCENT_YEL, 1.1)
    add_line(slide, 5.35, 5.72, 7.2, 5.72, ACCENT_RED, 1.1)
    return slide

# ─────────────────────────────────────────────────────────────────────────────
# 内容页布局
# ─────────────────────────────────────────────────────────────────────────────
def layout_three_cols(slide, cards_data, start_y=1.55, card_h=4.95):
    """三列大卡片。content 可为 list 或 dict(kpi/unit/desc/lines/tag/watermark)。"""
    card_w, gap = 3.72, 0.52
    for i, item in enumerate(cards_data[:3]):
        icon, title, c, content = item[0], item[1], item[2], item[3]
        x = 0.88 + i * (card_w + gap)
        add_rounded_rect(slide, x, start_y, card_w, card_h, BG_CARD, c, 0.9)
        add_rect(slide, x, start_y, card_w, 0.05, c)
        if isinstance(content, dict):
            kpi = content.get("kpi", "")
            unit = content.get("unit", "")
            desc = content.get("desc", "")
            lines = content.get("lines", [])
            tag = content.get("tag", "")
            watermark = content.get("watermark", "")
        else:
            kpi, unit, desc, lines, tag, watermark = "", "", "", list(content), "", ""
        if icon:
            add_circle(slide, x + card_w / 2 - 0.42, start_y + 0.36, 0.84, BG_CARD2, c, 0.8, transparency=0.15)
            add_text(slide, icon, x + card_w / 2 - 0.34, start_y + 0.50, 0.68, 0.36, 20, False, c, PP_ALIGN.CENTER)
        if watermark:
            add_text(slide, watermark, x + 0.65, start_y + 0.9, 2.4, 0.72, 42, False, RGBColor(0x46, 0x23, 0x10), PP_ALIGN.CENTER, FONT_TITLE)
        y = start_y + 1.35
        if kpi:
            add_text(slide, kpi, x + 0.15, y, card_w - 0.3, 0.70, 42, True, ACCENT_YEL, PP_ALIGN.CENTER)
            add_text(slide, unit, x + 0.15, y + 0.70, card_w - 0.3, 0.25, 12, False, ACCENT_YEL, PP_ALIGN.CENTER)
            y += 1.05
        add_text(slide, title, x + 0.18, y, card_w - 0.36, 0.34, 19, False, c if not kpi else TEXT_WHITE, PP_ALIGN.CENTER, FONT_TITLE)
        y += 0.43
        add_line(slide, x + 0.28, y, x + card_w - 0.28, y, BORDER_DIM, 0.6)
        y += 0.14
        if desc:
            add_text(slide, desc, x + 0.30, y, card_w - 0.60, 0.32, 11.5, False, TEXT_LIGHT, PP_ALIGN.CENTER)
            y += 0.36
        if lines:
            add_multiline_text(slide, lines, x + 0.35, y, card_w - 0.70, start_y + card_h - y - 0.65,
                               11.2, TEXT_LIGHT, PP_ALIGN.LEFT, False, FONT_BODY, 1.08)
        if tag:
            add_rounded_rect(slide, x + 0.45, start_y + card_h - 0.58, card_w - 0.9, 0.30, BG_DEEP, c, 0.6)
            add_text(slide, tag, x + 0.50, start_y + card_h - 0.52, card_w - 1.0, 0.18, 8.5, False, c, PP_ALIGN.CENTER)


def layout_kpi_2x2(slide, kpi_data, x=5.95, y=1.55, w=6.25, h=4.7):
    positions = [(x, y), (x + w / 2 + 0.18, y), (x, y + h / 2 + 0.18), (x + w / 2 + 0.18, y + h / 2 + 0.18)]
    card_w, card_h = w / 2 - 0.10, h / 2 - 0.10
    for (name, val, unit, color, compare, interpret, desc), (px, py) in zip(kpi_data[:4], positions):
        add_rounded_rect(slide, px, py, card_w, card_h, BG_CARD, color, 0.9)
        add_rect(slide, px, py, card_w, 0.05, color)
        add_text(slide, val, px + 0.18, py + 0.22, 1.4, 0.55, 29, True, color, PP_ALIGN.LEFT)
        add_text(slide, name, px + 1.45, py + 0.28, card_w - 1.62, 0.32, 15, False, TEXT_WHITE, font_name=FONT_TITLE)
        add_line(slide, px + 0.20, py + 0.88, px + card_w - 0.22, py + 0.88, BORDER_DIM, 0.6)
        lines = [s for s in [unit, compare, interpret, desc] if s]
        add_multiline_text(slide, lines, px + 0.20, py + 1.02, card_w - 0.4, card_h - 1.16, 10.2, TEXT_LIGHT)


def layout_timeline(slide, items, y=3.35):
    add_line(slide, 0.95, y, 12.35, y, BORDER_DIM, 1.2)
    n = len(items)
    for i, (year, title, color, desc) in enumerate(items):
        x = 1.05 + i * (11.0 / max(1, n - 1))
        add_circle(slide, x - 0.055, y - 0.055, 0.11, color, TEXT_WHITE, 0.8)
        up = i % 2 == 0
        if up:
            add_rounded_rect(slide, x - 0.45, y - 0.95, 0.90, 0.28, color, line_color=None)
            add_text(slide, year, x - 0.43, y - 0.90, 0.86, 0.16, 8.5, True, TEXT_WHITE, PP_ALIGN.CENTER)
            add_text(slide, title, x - 0.75, y - 0.55, 1.5, 0.28, 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
            add_text(slide, desc, x - 0.75, y - 0.30, 1.5, 0.22, 8, False, TEXT_DIM, PP_ALIGN.CENTER)
        else:
            add_rounded_rect(slide, x - 0.45, y + 0.65, 0.90, 0.28, color, line_color=None)
            add_text(slide, year, x - 0.43, y + 0.70, 0.86, 0.16, 8.5, True, BG_DARK if color == ACCENT_YEL else TEXT_WHITE, PP_ALIGN.CENTER)
            add_text(slide, title, x - 0.75, y + 0.28, 1.5, 0.28, 9, False, TEXT_LIGHT, PP_ALIGN.CENTER)
            add_text(slide, desc, x - 0.75, y + 0.50, 1.5, 0.20, 8, False, TEXT_DIM, PP_ALIGN.CENTER)


def layout_workflow(slide, stages, y=1.68):
    """流程页专用布局：上方 5 步小时间线 + 下方 3 个阶段卡。

    stages: [(step_label, title, color, subtitle, lines, tag), ...]，建议 3-5 项。
    目标是避免节点、Step 标签、三张卡片、底部 takeaway 互相重叠。
    """
    n = min(len(stages), 5)
    if n == 0:
        return
    line_y = y + 1.22
    add_line(slide, 1.05, line_y, 12.15, line_y, BORDER_DIM, 1.0)
    for i, (step, title, color, subtitle, lines, tag) in enumerate(stages[:n]):
        cx = 1.10 + i * (11.00 / max(1, n - 1))
        add_circle(slide, cx - 0.055, line_y - 0.055, 0.11, color, TEXT_WHITE, 0.7)
        if i % 2 == 0:
            label_y, title_y = line_y - 0.82, line_y - 0.50
        else:
            label_y, title_y = line_y + 0.34, line_y + 0.06
        add_rounded_rect(slide, cx - 0.42, label_y, 0.84, 0.25, color, line_color=None)
        add_text(slide, step, cx - 0.38, label_y + 0.055, 0.76, 0.12, 7.0, True,
                 BG_DARK if color == ACCENT_YEL else TEXT_WHITE, PP_ALIGN.CENTER)
        add_text(slide, title, cx - 0.70, title_y, 1.40, 0.18, 8.0, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    cards = stages[:3]
    card_w, gap, card_y, card_h = 3.55, 0.72, y + 2.35, 1.72
    for i, (step, title, color, subtitle, lines, tag) in enumerate(cards):
        x = 1.05 + i * (card_w + gap)
        add_rounded_rect(slide, x, card_y, card_w, card_h, BG_CARD, color, 1.0)
        add_rect(slide, x, card_y, card_w, 0.055, color)
        add_circle(slide, x + card_w / 2 - 0.31, card_y + 0.28, 0.62, BG_CARD2, color, 0.8)
        add_text(slide, str(i + 1), x + card_w / 2 - 0.18, card_y + 0.43, 0.36, 0.18, 16, False, color, PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.28, card_y + 0.98, card_w - 0.56, 0.22, 13.2, False, color, PP_ALIGN.CENTER, FONT_TITLE)
        if tag:
            add_rounded_rect(slide, x + 0.42, card_y + 1.25, card_w - 0.84, 0.22, BG_DEEP, color, 0.7)
            add_text(slide, tag, x + 0.48, card_y + 1.295, card_w - 0.96, 0.10, 7.2, False, color, PP_ALIGN.CENTER)
        if lines:
            add_multiline_text(slide, lines[:2], x + 0.34, card_y + 1.48, card_w - 0.68, 0.35, 8.2, TEXT_LIGHT, PP_ALIGN.LEFT)


def layout_value_split(slide, left_title, left_items, summary_title, summary_lines,
                       audience_blocks=None, y=2.04, takeaway="", takeaway_y=6.32):
    """应用场景/产品价值页：左侧要点栈 + 右侧总结卡 + 右下受众卡。

    适合"应用场景与产品价值""核心亮点"等页面。所有文本和卡片均保留可编辑。
    left_items: [(title, desc, color), ...]
    audience_blocks: [(title, desc, color), ...]
    """
    left_x, left_w, panel_h = 0.78, 5.55, 4.18
    right_x, right_w = 6.70, 5.75
    add_rounded_rect(slide, left_x, y, left_w, panel_h, BG_CARD, ACCENT_RED, 1.0)
    add_rect(slide, left_x, y, left_w, 0.055, ACCENT_RED)
    add_text(slide, f"· {left_title} ·", left_x + 1.65, y + 0.22, 2.4, 0.22, 12.2, False, ACCENT_YEL, PP_ALIGN.CENTER, FONT_TITLE)
    row_h = min(0.72, (panel_h - 0.88) / max(1, len(left_items[:4])))
    for i, (name, desc, color) in enumerate(left_items[:4]):
        iy = y + 0.62 + i * row_h
        add_rounded_rect(slide, left_x + 0.24, iy, left_w - 0.48, row_h - 0.12, BG_DEEP, line_color=None)
        add_rect(slide, left_x + 0.24, iy, 0.04, row_h - 0.12, color)
        add_text(slide, name, left_x + 0.45, iy + 0.13, 1.15, 0.18, 8.6, True, TEXT_WHITE)
        add_text(slide, desc, left_x + 1.58, iy + 0.13, left_w - 2.05, 0.20, 8.2, False, TEXT_DIM)
    add_rounded_rect(slide, right_x, y, right_w, 1.62, BG_CARD, ACCENT_YEL, 1.0)
    add_text(slide, summary_title, right_x + 0.35, y + 0.40, right_w - 0.70, 0.24, 13.2, False, ACCENT_YEL, PP_ALIGN.CENTER, FONT_TITLE)
    add_multiline_text(slide, summary_lines[:4], right_x + 0.58, y + 0.77, right_w - 1.16, 0.58, 8.8, TEXT_LIGHT, PP_ALIGN.LEFT)
    audience_blocks = audience_blocks or []
    block_y = y + 2.80
    for i, (title, desc, color) in enumerate(audience_blocks[:2]):
        by = block_y + i * 0.86
        add_rounded_rect(slide, right_x, by, right_w, 0.68, BG_CARD, color, 0.9)
        add_text(slide, title, right_x + 0.35, by + 0.14, right_w - 0.70, 0.18, 11.0, True, color, PP_ALIGN.CENTER)
        add_text(slide, desc, right_x + 0.42, by + 0.39, right_w - 0.84, 0.16, 8.2, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    if takeaway:
        add_imperial_takeaway(slide, takeaway, y=takeaway_y, color=ACCENT_YEL)


def layout_bullet_list(slide, items, start_y=1.62, x=0.86, w=11.6, row_h=0.78):
    for i, item in enumerate(items[:6]):
        icon, title, color, desc = item if len(item) == 4 else ("", item[0], item[1], item[2])
        y = start_y + i * row_h
        add_rounded_rect(slide, x, y, w, row_h - 0.10, BG_CARD, line_color=None, transparency=0.04)
        add_rect(slide, x, y, 0.045, row_h - 0.10, color)
        add_circle(slide, x + 0.22, y + 0.23, 0.22, BG_CARD2, color, 0.7)
        add_text(slide, icon or str(i + 1), x + 0.22, y + 0.275, 0.22, 0.10, 7.5, False, color, PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.55, y + 0.10, w - 0.7, 0.25, 15, False, color, font_name=FONT_TITLE)
        add_text(slide, desc, x + 0.55, y + 0.38, w - 0.7, 0.22, 10.5, False, TEXT_LIGHT)


def layout_quote(slide, quote_lines, source=""):
    add_text(slide, "\u201c", 0.95, 2.0, 1.0, 0.8, 68, False, RGBColor(0x58, 0x36, 0x16), font_name="Georgia")
    add_multiline_text(slide, quote_lines, 1.55, 2.72, 10.8, 1.2, 26, TEXT_WHITE, PP_ALIGN.LEFT, False, FONT_TITLE, 1.12)
    if source:
        add_text(slide, f"— {source}", 9.25, 4.35, 2.8, 0.28, 11, False, TEXT_LIGHT, PP_ALIGN.RIGHT)


def layout_table(slide, headers, rows, col_widths=None, x0=0.82, y0=1.62, table_w=11.7, bottom_y=5.95):
    if not col_widths:
        col_widths = [table_w / len(headers)] * len(headers)
    row_h = min(0.58, max(0.42, (bottom_y - y0) / (len(rows) + 1)))
    xs = [x0]
    for cw in col_widths[:-1]:
        xs.append(xs[-1] + cw)
    add_rounded_rect(slide, x0, y0, sum(col_widths), row_h, BG_CARD2, ACCENT_YEL, 0.8)
    for j, h in enumerate(headers):
        add_text(slide, h, xs[j] + 0.08, y0 + row_h * 0.30, col_widths[j] - 0.16, row_h * 0.35,
                 10.5, True, ACCENT_YEL, PP_ALIGN.CENTER)
    for i, row in enumerate(rows[:10]):
        y = y0 + (i + 1) * row_h
        add_rect(slide, x0, y, sum(col_widths), row_h, BG_CARD if i % 2 == 0 else BG_DEEP, BORDER_SOFT, 0.2)
        for j, cell in enumerate(row):
            add_text(slide, cell, xs[j] + 0.08, y + row_h * 0.32, col_widths[j] - 0.16, row_h * 0.32,
                     9.5, False, TEXT_LIGHT, PP_ALIGN.CENTER)


def layout_metrics_strip(slide, metrics, x=0.68, y=1.52, w=12.0, h=0.50):
    """顶部指标条：适合 3-5 个短指标。自动压缩长标签字号，避免挤压。"""
    add_rounded_rect(slide, x, y, w, h, BG_CARD, BORDER_DIM, 0.8)
    n = len(metrics)
    cell_w = w / n
    for i, (value, label, color) in enumerate(metrics):
        cx = x + i * cell_w
        if i:
            add_line(slide, cx, y + 0.08, cx, y + h - 0.08, BORDER_DIM, 0.6)
        add_text(slide, value, cx + 0.05, y + 0.07, cell_w - 0.10, 0.22, 15.5, True, color, PP_ALIGN.CENTER)
        add_text(slide, label, cx + 0.05, y + 0.30, cell_w - 0.10, 0.15, 7.5, False, TEXT_LIGHT, PP_ALIGN.CENTER)


def layout_two_panels(slide, left_title, left_items, right_blocks, y=2.25):
    add_rounded_rect(slide, 0.68, y, 5.65, 4.35, BG_CARD, ACCENT_RED, 1.0)
    add_rect(slide, 0.68, y, 5.65, 0.055, ACCENT_RED)
    add_text(slide, f"· {left_title} ·", 1.9, y + 0.22, 2.6, 0.26, 14, False, ACCENT_YEL, PP_ALIGN.CENTER, FONT_TITLE)
    for i, (name, desc, color) in enumerate(left_items[:4]):
        iy = y + 0.60 + i * 0.82
        add_rounded_rect(slide, 0.92, iy, 4.9, 0.62, BG_DEEP, line_color=None)
        add_rect(slide, 0.92, iy, 0.04, 0.62, color)
        add_text(slide, name, 1.12, iy + 0.12, 1.4, 0.20, 13, False, TEXT_WHITE, font_name=FONT_TITLE)
        add_text(slide, desc, 2.10, iy + 0.13, 3.45, 0.20, 8.8, False, TEXT_DIM)
    add_rounded_rect(slide, 6.70, y, 5.85, 1.70, BG_CARD, ACCENT_YEL, 1.0)
    add_rect(slide, 6.70, y, 5.85, 0.055, ACCENT_YEL)
    for i, block in enumerate(right_blocks[:3]):
        bx, by, bw, bh = (6.70, y + 1.95 + i * 0.88, 5.85, 0.70) if i else (6.95, y + 0.25, 5.35, 1.15)
        if i:
            add_rounded_rect(slide, bx, by, bw, bh, BG_CARD, block.get("color", ACCENT_YEL), 0.8)
        add_text(slide, block.get("title", ""), bx + 0.25, by + 0.14, bw - 0.5, 0.24, 14, False, block.get("color", ACCENT_YEL), PP_ALIGN.CENTER, FONT_TITLE)
        add_multiline_text(slide, block.get("lines", []), bx + 0.35, by + 0.42, bw - 0.7, bh - 0.48, 9.5, TEXT_LIGHT, PP_ALIGN.LEFT)


def layout_image_left_grid(slide, image_path, caption, cards, image_box=(0.68, 1.55, 4.95, 4.95)):
    x, y, w, h = image_box
    add_rounded_rect(slide, x, y, w, h, BG_CARD, BORDER_DIM, 0.8)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        # 可编辑通用占位插画：抽象场景/架构/路径，不绑定任何固定行业或历史题材。
        add_rect(slide, x + 0.05, y + 0.05, w - 0.10, h - 0.10, RGBColor(0x5E, 0x45, 0x25), transparency=0.15)
        add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + 0.4, y + 1.2, 1.4, 1.2, RGBColor(0x3C, 0x4B, 0x38), None, 0)
        add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + 1.1, y + 0.9, 1.8, 1.5, RGBColor(0x4B, 0x5A, 0x44), None, 0)
        for k in range(7):
            add_line(slide, x + 0.35 + k * 0.55, y + 3.25, x + 0.10 + k * 0.55, y + 4.55, ACCENT_YEL, 1.0)
        for k in range(5):
            add_circle(slide, x + 2.2 + k * 0.42, y + 3.5 + (k % 2) * 0.18, 0.18, ACCENT_ORANGE, None)
    add_rounded_rect(slide, x + 0.20, y + h - 0.48, w - 0.40, 0.28, BG_INK, line_color=None, transparency=0.15)
    add_text(slide, caption, x + 0.25, y + h - 0.42, w - 0.50, 0.16, 8.5, False, ACCENT_GRN, PP_ALIGN.CENTER)
    layout_kpi_2x2(slide, cards, x=5.95, y=1.55, w=6.25, h=4.7)


def layout_art_showcase(slide, left_items, center_title, center_subtitle, features):
    """左侧样本列表 + 中央可编辑抽象视觉 + 右侧特征列表，适合理念、品牌、案例、产品体验页。"""
    for i, (title, desc, color) in enumerate(left_items[:3]):
        y = 1.58 + i * 1.55
        add_rounded_rect(slide, 0.68, y, 4.05, 1.25, BG_CARD, BORDER_DIM, 0.8)
        add_rect(slide, 0.68, y, 0.045, 1.25, color)
        add_circle(slide, 0.95, y + 0.35, 0.55, BG_CARD2, color, 0.8)
        add_text(slide, title, 1.78, y + 0.28, 1.4, 0.26, 15, False, TEXT_WHITE, font_name=FONT_TITLE)
        add_multiline_text(slide, desc, 1.78, y + 0.60, 2.6, 0.45, 8.5, TEXT_DIM)
    add_rounded_rect(slide, 5.05, 1.58, 7.50, 5.00, BG_CARD, ACCENT_YEL, 1.0)
    add_rect(slide, 5.05, 1.58, 7.50, 0.055, ACCENT_YEL)
    add_text(slide, f"\u00b7 {center_title} \u00b7", 6.95, 1.85, 3.8, 0.28, 15, False, ACCENT_YEL, PP_ALIGN.CENTER, FONT_TITLE)
    # 抽象视觉图：全部由原生椭圆/矩形组合，非图片
    add_rounded_rect(slide, 5.38, 2.30, 4.75, 3.08, BG_DEEP, ACCENT_YEL, 0.8)
    add_circle(slide, 7.30, 2.72, 0.58, ACCENT_CREAM, BORDER_DIM, 0.4)
    add_shape(slide, MSO_SHAPE.OVAL, 7.10, 3.26, 0.70, 1.55, ACCENT_RED, BORDER_DIM, 0.5)
    add_shape(slide, MSO_SHAPE.OVAL, 6.35, 3.35, 0.80, 0.36, ACCENT_YEL, BORDER_DIM, 0.5)
    add_shape(slide, MSO_SHAPE.OVAL, 7.70, 3.15, 0.72, 0.34, ACCENT_YEL, BORDER_DIM, 0.5)
    add_circle(slide, 8.55, 2.80, 0.50, RGBColor(0x66, 0x58, 0x46), BORDER_DIM, 0.5)
    add_shape(slide, MSO_SHAPE.OVAL, 8.35, 3.28, 0.50, 1.36, RGBColor(0x29, 0x45, 0x32), BORDER_DIM, 0.5)
    add_text(slide, center_subtitle, 5.65, 5.45, 4.3, 0.18, 8.5, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    add_rounded_rect(slide, 10.45, 2.30, 1.85, 3.08, BG_DEEP, BORDER_DIM, 0.8)
    add_text(slide, "关键特征", 10.72, 2.55, 1.3, 0.24, 11, False, ACCENT_YEL, PP_ALIGN.CENTER, FONT_TITLE)
    for i, (label, color) in enumerate(features[:5]):
        fy = 2.95 + i * 0.42
        add_circle(slide, 10.62, fy + 0.05, 0.08, color, color)
        add_text(slide, label, 10.82, fy, 1.2, 0.18, 8.5, False, TEXT_LIGHT)


def layout_cards_grid(slide, cards, cols=2, x=0.78, y=1.55, w=11.85, h=4.95, gap_x=0.36, gap_y=0.28):
    """通用多卡片网格：适合 4–6 个信息点、能力点、风险点、案例要点。cards 为 dict 列表。"""
    if not cards:
        return
    cols = max(1, min(cols, 3))
    rows = (min(len(cards), 6) + cols - 1) // cols
    card_w = (w - gap_x * (cols - 1)) / cols
    card_h = (h - gap_y * (rows - 1)) / rows
    for i, card in enumerate(cards[:6]):
        col, row = i % cols, i // cols
        cx, cy = x + col * (card_w + gap_x), y + row * (card_h + gap_y)
        color = card.get("color", BRAND_COLORS[i % len(BRAND_COLORS)])
        add_rounded_rect(slide, cx, cy, card_w, card_h, BG_CARD, color, 0.85)
        add_rect(slide, cx, cy, 0.05, card_h, color)
        tag = card.get("tag", f"{i + 1:02d}")
        add_rounded_rect(slide, cx + 0.20, cy + 0.18, 0.72, 0.25, color, line_color=None)
        add_text(slide, tag, cx + 0.24, cy + 0.23, 0.64, 0.12, 7.5, True,
                 BG_DARK if color == ACCENT_YEL else TEXT_WHITE, PP_ALIGN.CENTER)
        add_text(slide, card.get("title", ""), cx + 1.02, cy + 0.18, card_w - 1.25, 0.26,
                 13.8, False, TEXT_WHITE, font_name=FONT_TITLE)
        if card.get("desc"):
            add_text(slide, card.get("desc", ""), cx + 0.24, cy + 0.58, card_w - 0.48, 0.22,
                     9.2, False, TEXT_LIGHT)
        lines = card.get("lines", []) or []
        if lines:
            add_multiline_text(slide, lines[:3], cx + 0.28, cy + 0.88, card_w - 0.56,
                               max(0.25, card_h - 1.06), 8.8, TEXT_DIM)


def layout_focus_statement(slide, statement, evidences=None, y=1.72, color=ACCENT_YEL):
    """重点观点页：中央大判断 + 下方证据卡，适合观点阐述、模块结论、管理层摘要。"""
    add_rounded_rect(slide, 1.05, y, 11.25, 1.60, BG_CARD, color, 1.0)
    add_text(slide, statement, 1.55, y + 0.42, 10.25, 0.62, 25, False, TEXT_WHITE, PP_ALIGN.CENTER, FONT_TITLE)
    add_line(slide, 3.0, y + 1.22, 6.05, y + 1.22, color, 0.9)
    add_line(slide, 7.25, y + 1.22, 10.35, y + 1.22, color, 0.9)
    add_text(slide, "\u25c6", 6.44, y + 1.10, 0.45, 0.22, 12, False, color, PP_ALIGN.CENTER)
    evidences = evidences or []
    card_w = 3.55
    for i, item in enumerate(evidences[:3]):
        title, desc, c = item if len(item) >= 3 else (item[0], item[1], BRAND_COLORS[i % 4])
        x = 1.08 + i * (card_w + 0.64)
        cy = y + 2.08
        add_rounded_rect(slide, x, cy, card_w, 1.65, BG_CARD, c, 0.85)
        add_rect(slide, x, cy, card_w, 0.05, c)
        add_text(slide, title, x + 0.25, cy + 0.28, card_w - 0.5, 0.24, 13, False, c, PP_ALIGN.CENTER, FONT_TITLE)
        add_text(slide, desc, x + 0.30, cy + 0.70, card_w - 0.60, 0.55, 9.2, False, TEXT_LIGHT, PP_ALIGN.CENTER)


def layout_comparison(slide, left, right, conclusion="", y=1.62):
    """左右对比页：适合方案对比、前后变化、机会/风险、竞品差异。left/right 为 dict。"""
    panels = [(0.78, left, left.get("color", ACCENT_RED)), (7.00, right, right.get("color", ACCENT_GRN))]
    for x, data, color in panels:
        add_rounded_rect(slide, x, y, 5.55, 4.55, BG_CARD, color, 1.0)
        add_rect(slide, x, y, 5.55, 0.06, color)
        add_text(slide, data.get("title", ""), x + 0.35, y + 0.32, 4.85, 0.32, 18, False, color, PP_ALIGN.CENTER, FONT_TITLE)
        if data.get("subtitle"):
            add_text(slide, data.get("subtitle", ""), x + 0.45, y + 0.78, 4.65, 0.20, 9.5, False, TEXT_LIGHT, PP_ALIGN.CENTER)
        add_line(slide, x + 0.42, y + 1.12, x + 5.10, y + 1.12, BORDER_DIM, 0.7)
        points = data.get("points", []) or []
        for i, point in enumerate(points[:5]):
            py = y + 1.45 + i * 0.50
            add_circle(slide, x + 0.50, py + 0.06, 0.08, color, color)
            add_text(slide, point, x + 0.72, py, 4.35, 0.20, 9.5, False, TEXT_LIGHT)
    add_circle(slide, 6.34, y + 2.12, 0.44, BG_CARD2, ACCENT_YEL, 0.8)
    add_text(slide, "VS", 6.40, y + 2.28, 0.32, 0.12, 8.5, True, ACCENT_YEL, PP_ALIGN.CENTER)
    if conclusion:
        add_imperial_takeaway(slide, conclusion, y=6.28, color=ACCENT_YEL)


def layout_case_study(slide, case_title, context, actions, results, takeaway="", y=1.58):
    """通用案例页：场景/挑战—关键动作—结果/启示，适合案例、复盘、最佳实践。"""
    add_rounded_rect(slide, 0.78, y, 4.05, 4.85, BG_CARD, ACCENT_RED, 1.0)
    add_rect(slide, 0.78, y, 4.05, 0.06, ACCENT_RED)
    add_text(slide, case_title, 1.05, y + 0.30, 3.50, 0.32, 17, False, ACCENT_YEL, PP_ALIGN.CENTER, FONT_TITLE)
    add_multiline_text(slide, context[:5], 1.05, y + 0.92, 3.50, 2.6, 10, TEXT_LIGHT)
    add_rounded_rect(slide, 1.08, y + 3.85, 3.45, 0.36, BG_DEEP, ACCENT_RED, 0.7)
    add_text(slide, "场景 \u00b7 问题 \u00b7 约束", 1.20, y + 3.94, 3.20, 0.12, 8, False, ACCENT_RED, PP_ALIGN.CENTER)
    sections = [("关键动作", actions, ACCENT_YEL, 5.18), ("结果与经验", results, ACCENT_GRN, 9.08)]
    for title, lines, color, x in sections:
        add_rounded_rect(slide, x, y, 3.45, 4.85, BG_CARD, color, 1.0)
        add_rect(slide, x, y, 3.45, 0.06, color)
        add_text(slide, title, x + 0.30, y + 0.35, 2.85, 0.28, 16, False, color, PP_ALIGN.CENTER, FONT_TITLE)
        for i, line in enumerate(lines[:5]):
            py = y + 0.95 + i * 0.56
            add_circle(slide, x + 0.35, py + 0.06, 0.08, color, color)
            add_text(slide, line, x + 0.55, py, 2.55, 0.22, 9, False, TEXT_LIGHT)
    if takeaway:
        add_imperial_takeaway(slide, takeaway, y=6.32, color=ACCENT_YEL)


def layout_conclusion(slide, main_points, next_steps=None, y=1.58):
    """总结页：左侧结论清单 + 右侧下一步行动，适合全篇总结或模块收束。"""
    add_rounded_rect(slide, 0.78, y, 7.10, 4.90, BG_CARD, ACCENT_YEL, 1.0)
    add_rect(slide, 0.78, y, 7.10, 0.06, ACCENT_YEL)
    add_text(slide, "核心结论", 3.0, y + 0.28, 2.6, 0.30, 18, False, ACCENT_YEL, PP_ALIGN.CENTER, FONT_TITLE)
    for i, point in enumerate(main_points[:5]):
        py = y + 0.92 + i * 0.68
        color = BRAND_COLORS[i % len(BRAND_COLORS)]
        add_circle(slide, 1.10, py + 0.02, 0.26, BG_CARD2, color, 0.7)
        add_text(slide, str(i + 1), 1.17, py + 0.08, 0.12, 0.08, 7.5, True, color, PP_ALIGN.CENTER)
        add_text(slide, point, 1.55, py, 5.85, 0.24, 10.8, False, TEXT_LIGHT)
    next_steps = next_steps or []
    add_rounded_rect(slide, 8.25, y, 4.10, 4.90, BG_CARD, ACCENT_GRN, 1.0)
    add_rect(slide, 8.25, y, 4.10, 0.06, ACCENT_GRN)
    add_text(slide, "下一步", 9.15, y + 0.28, 2.3, 0.30, 18, False, ACCENT_GRN, PP_ALIGN.CENTER, FONT_TITLE)
    for i, step in enumerate(next_steps[:4]):
        py = y + 0.95 + i * 0.78
        add_rounded_rect(slide, 8.58, py, 3.45, 0.52, BG_DEEP, line_color=None)
        add_rect(slide, 8.58, py, 0.04, 0.52, ACCENT_GRN)
        add_text(slide, step, 8.78, py + 0.15, 3.02, 0.16, 8.8, False, TEXT_LIGHT)


def layout_badge_list(slide, badge_title, badge_subtitle, items, badge_meta="", y=1.58, color=ACCENT_YEL):
    """左侧圆形徽章 + 右侧堆叠信息条：适合核心概念解释、关键结论、模块总结、转折页。"""
    add_circle(slide, 1.05, y + 0.58, 2.55, BG_CARD, BORDER_DIM, 0.9, transparency=0.06)
    add_circle(slide, 1.30, y + 0.83, 2.05, BG_DEEP, color, 0.7, transparency=0.10)
    title_size = _fit_font_size(badge_title, 1.75, 34, min_size=22, density=1.8)
    add_text(slide, badge_title, 1.48, y + 1.48, 1.70, 0.42, title_size, False, color, PP_ALIGN.CENTER, FONT_TITLE, auto_fit=False)
    if badge_subtitle:
        add_text(slide, badge_subtitle, 1.45, y + 1.98, 1.75, 0.22, 9.0, False, TEXT_LIGHT, PP_ALIGN.CENTER)
    if badge_meta:
        add_rounded_rect(slide, 1.58, y + 2.42, 1.45, 0.32, BG_CARD2, color, 0.7)
        add_text(slide, badge_meta, 1.65, y + 2.50, 1.30, 0.12, 7.6, False, color, PP_ALIGN.CENTER)
    add_line(slide, 4.25, y + 0.15, 4.25, y + 4.62, BORDER_DIM, 0.8)
    for i, (title, desc, item_color) in enumerate(items[:6]):
        row_y = y + 0.18 + i * 0.72
        add_rounded_rect(slide, 4.55, row_y, 7.70, 0.55, BG_CARD, line_color=None, transparency=0.04)
        add_rect(slide, 4.55, row_y, 0.045, 0.55, item_color)
        add_circle(slide, 4.82, row_y + 0.17, 0.20, BG_CARD2, item_color, 0.7)
        add_text(slide, str(i + 1), 4.88, row_y + 0.225, 0.08, 0.06, 6.8, True, item_color, PP_ALIGN.CENTER)
        add_text(slide, title, 5.16, row_y + 0.09, 2.15, 0.18, 10.8, True, item_color)
        add_text(slide, desc, 7.20, row_y + 0.10, 4.78, 0.18, 8.8, False, TEXT_LIGHT)


def save_presentation(prs, export_dir: str, filename_stem: str) -> str:
    os.makedirs(export_dir, exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = os.path.join(export_dir, f"{filename_stem}_{ts}.pptx")
    prs.save(out)
    print(f"PPTX saved: {out}")
    return out
